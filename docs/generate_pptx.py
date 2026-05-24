from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color palette ──
C_BG       = RGBColor(0x0F, 0x17, 0x2A)   # dark navy
C_ACCENT   = RGBColor(0x3B, 0x82, 0xF6)   # blue
C_DANGER   = RGBColor(0xEF, 0x44, 0x44)   # red
C_WARN     = RGBColor(0xF5, 0x9E, 0x0B)   # amber
C_SUCCESS  = RGBColor(0x10, 0xB9, 0x81)   # green
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED    = RGBColor(0x94, 0xA3, 0xB8)
C_CARD     = RGBColor(0x1E, 0x29, 0x3B)
C_BORDER   = RGBColor(0x33, 0x44, 0x55)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]  # totally blank

# ─────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color if color else C_WHITE
    return tb

def bg(slide, color=C_BG):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=color)

def accent_bar(slide, y=0.55, h=0.045):
    add_rect(slide, 0.5, y, 1.6, h, fill=C_ACCENT)

def slide_label(slide, text, y=0.6):
    add_text(slide, text, 2.25, y-0.12, 10, 0.35, size=9,
             color=C_MUTED, bold=False)

def title_block(slide, title, subtitle=None, title_y=0.85, sub_y=1.45):
    add_text(slide, title, 0.5, title_y, 12.3, 0.7,
             size=36, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.5, sub_y, 12.3, 0.45,
                 size=16, color=C_MUTED, align=PP_ALIGN.LEFT)

def card(slide, x, y, w, h, title, body_lines, title_color=C_ACCENT):
    add_rect(slide, x, y, w, h, fill=C_CARD, line=C_BORDER, line_w=0.5)
    add_text(slide, title, x+0.18, y+0.12, w-0.3, 0.32,
             size=11, bold=True, color=title_color)
    body = "\n".join(body_lines)
    add_text(slide, body, x+0.18, y+0.44, w-0.3, h-0.55,
             size=9.5, color=C_MUTED)

def divider(slide, y):
    add_rect(slide, 0.5, y, 12.33, 0.015, fill=C_BORDER)

# ─────────────────────────────────────────
# SLIDE 1 — Cover
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)
# gradient strip
add_rect(s, 0, 0, 13.33, 7.5, fill=C_BG)
add_rect(s, 0, 5.2, 13.33, 2.3, fill=RGBColor(0x1A, 0x27, 0x3D))
add_rect(s, 0, 0, 0.08, 7.5, fill=C_ACCENT)

add_text(s, "SMARTGUARD AI", 0.5, 1.4, 12, 1.1, size=54, bold=True, color=C_WHITE)
add_text(s, "Smart Warehouse Bio-Hazard & Pest Detection System",
         0.5, 2.55, 12, 0.55, size=20, color=C_ACCENT, bold=True)
add_rect(s, 0.5, 3.3, 3.5, 0.045, fill=C_ACCENT)
add_text(s, "AI Open Innovation Challenge 2026  ·  PT. Kawan Lama Surveillance",
         0.5, 3.5, 12, 0.4, size=13, color=C_MUTED)
add_text(s, "Group 5  ·  Risly  ·  Sultan  ·  Fathir  ·  Misha",
         0.5, 4.05, 12, 0.35, size=12, color=C_MUTED)

# tech pills
for i, (label, col) in enumerate([("YOLO11", C_ACCENT), ("FastAPI", C_SUCCESS),
                                    ("React", C_ACCENT), ("SQLite", C_WARN), ("WebSocket", C_DANGER)]):
    add_rect(s, 0.5 + i*2.35, 5.6, 2.0, 0.42, fill=C_CARD, line=col, line_w=0.8)
    add_text(s, label, 0.5 + i*2.35, 5.62, 2.0, 0.38, size=11.5,
             bold=True, color=col, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# SLIDE 2 — Problem Statement
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)
slide_label(s, "PROBLEM STATEMENT")
title_block(s, "The Challenge at PT. Kawan Lama",
            "Large-scale warehouses face invisible threats that manual monitoring cannot catch in time.")

problems = [
    ("Financial Loss", "Rodents, snakes, and pests contaminate stored goods\ncausing millions in product damage yearly."),
    ("Worker Safety", "Venomous snakes (kobra, welang) have been found in\nwarehouse zones with no early warning system."),
    ("Manual Limits", "Guards cannot watch 200+ CCTV cameras simultaneously.\nThreats go undetected for hours."),
    ("No Audit Trail", "When incidents happen, there is no timestamped record\nof what happened, where, and when."),
]
for i, (t, b) in enumerate(problems):
    col = i % 2
    row = i // 2
    card(s, 0.5 + col*6.45, 2.3 + row*2.15, 6.15, 2.0, t, b.split("\n"),
         title_color=C_DANGER)

# ─────────────────────────────────────────
# SLIDE 3 — Our Solution
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)
slide_label(s, "OUR SOLUTION")
title_block(s, "SmartGuard AI — Automated Surveillance",
            "A full-stack AI system that watches every zone 24/7 and alerts the right people instantly.")

features = [
    ("Real-Time Detection", C_ACCENT,
     ["Custom YOLO11 model", "Snake / Cat / Gecko / Lizard", "Runs on CPU, no GPU needed"]),
    ("HUD Bounding Boxes", C_SUCCESS,
     ["Color-coded by risk level", "Confidence score overlay", "Annotated snapshots saved"]),
    ("Instant Alerts", C_DANGER,
     ["WebSocket push in < 1 sec", "Share to WhatsApp / Telegram", "Toast + audio TTS warning"]),
    ("Multi-Zone Monitoring", C_WARN,
     ["4 parallel camera zones", "Independent inference threads", "Per-zone start/stop control"]),
    ("Analytics Dashboard", C_ACCENT,
     ["Weekly trend charts", "Zone activity heatmap", "Risk level distribution"]),
    ("Secure Access", C_SUCCESS,
     ["Role-based access control", "Invite-only registration", "bcrypt + rate limiting"]),
]
cols = 3
for i, (title, col, lines) in enumerate(features):
    c = i % cols
    r = i // cols
    card(s, 0.5 + c*4.28, 2.3 + r*2.1, 4.08, 2.0, title, lines, title_color=col)

# ─────────────────────────────────────────
# SLIDE 4 — System Architecture
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)
slide_label(s, "SYSTEM ARCHITECTURE")
title_block(s, "How It Works — End to End Pipeline", sub_y=1.42)

# Pipeline boxes
steps = [
    ("Camera\n/ Video", C_MUTED),
    ("OpenCV\nCapture", C_ACCENT),
    ("YOLO11\nInference", C_DANGER),
    ("FastAPI\nBackend", C_SUCCESS),
    ("SQLite\nDatabase", C_WARN),
    ("WebSocket\nBroadcast", C_ACCENT),
    ("React\nDashboard", C_SUCCESS),
]
box_w = 1.55
gap   = 0.25
start_x = 0.35
y_box = 2.5

for i, (label, col) in enumerate(steps):
    x = start_x + i*(box_w + gap)
    add_rect(s, x, y_box, box_w, 1.1, fill=C_CARD, line=col, line_w=1.2)
    add_text(s, label, x, y_box, box_w, 1.1, size=11, bold=True,
             color=col, align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        ax = x + box_w + 0.02
        add_text(s, "→", ax, y_box+0.3, gap+0.02, 0.5,
                 size=16, bold=True, color=C_MUTED, align=PP_ALIGN.CENTER)

# Feature highlights below
highlights = [
    ("< 200ms", "Per frame\ninference"),
    ("4 Zones", "Parallel\nthreads"),
    ("< 1 sec", "Alert\ndelivery"),
    ("100%", "Real data\nno mocks"),
]
for i, (val, lbl) in enumerate(highlights):
    x = 1.0 + i*3.1
    add_text(s, val, x, 4.2, 3.0, 0.6, size=28, bold=True,
             color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, lbl, x, 4.8, 3.0, 0.45, size=11,
             color=C_MUTED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# SLIDE 5 — AI Model Performance
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)
slide_label(s, "AI MODEL PERFORMANCE")
title_block(s, "YOLO11 Nano — Detection Accuracy",
            "Trained on 2,000+ annotated images with data augmentation for low-light warehouse conditions.")

metrics = [
    ("mAP@0.5",   "71%",  C_SUCCESS),
    ("Precision",  "74%",  C_ACCENT),
    ("Recall",     "68%",  C_WARN),
    ("Inference",  "~180ms", C_DANGER),
]
for i, (label, val, col) in enumerate(metrics):
    x = 0.5 + i*3.2
    add_rect(s, x, 2.3, 3.0, 1.8, fill=C_CARD, line=col, line_w=1.2)
    add_text(s, val, x, 2.5, 3.0, 0.9, size=38, bold=True,
             color=col, align=PP_ALIGN.CENTER)
    add_text(s, label, x, 3.35, 3.0, 0.4, size=12,
             color=C_MUTED, align=PP_ALIGN.CENTER)

classes = [
    ("Snake",  "Bio-Hazard",  C_DANGER,  "65%+ confidence on\ntest footage"),
    ("Cat",    "Contamination", C_WARN,  "72%+ confidence,\nhighest accuracy class"),
    ("Gecko",  "Monitoring",  C_SUCCESS, "63%+ after low-light\naugmentation"),
    ("Lizard", "Monitoring",  C_SUCCESS, "61%+ confidence on\ntest footage"),
]
for i, (name, risk, col, note) in enumerate(classes):
    x = 0.5 + i*3.2
    add_rect(s, x, 4.3, 3.0, 2.1, fill=C_CARD, line=C_BORDER, line_w=0.5)
    add_text(s, name, x+0.15, 4.42, 2.7, 0.38, size=14, bold=True, color=C_WHITE)
    add_rect(s, x+0.15, 4.82, 1.6, 0.28, fill=col)
    add_text(s, risk, x+0.15, 4.82, 1.6, 0.28, size=8.5,
             bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, note, x+0.15, 5.18, 2.7, 0.85, size=9, color=C_MUTED)

# ─────────────────────────────────────────
# SLIDE 6 — Tech Stack
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)
slide_label(s, "TECHNOLOGY STACK")
title_block(s, "Built With Production-Grade Tools",
            "Every component chosen for reliability, speed, and ease of deployment.")

stack = [
    ("Frontend",  C_ACCENT, [
        "React 19 + Vite",
        "Recharts (analytics)",
        "React Router v6",
        "CSS custom design system",
        "WebSocket client"
    ]),
    ("Backend",   C_SUCCESS, [
        "FastAPI (Python)",
        "Uvicorn ASGI server",
        "OpenCV (video capture)",
        "python-jose (JWT)",
        "bcrypt (passwords)"
    ]),
    ("AI / ML",   C_DANGER, [
        "YOLO11 Nano model",
        "Ultralytics framework",
        "Custom dataset 2000+ images",
        "COCO pretrained base",
        "CPU inference ready"
    ]),
    ("Infra",     C_WARN, [
        "SQLite (persistent DB)",
        "WebSocket manager",
        "MJPEG video streaming",
        "Docker + docker-compose",
        "GitHub (version control)"
    ]),
]
for i, (title, col, items) in enumerate(stack):
    x = 0.5 + i*3.2
    add_rect(s, x, 2.3, 3.0, 4.5, fill=C_CARD, line=col, line_w=1.2)
    add_rect(s, x, 2.3, 3.0, 0.45, fill=col)
    add_text(s, title, x, 2.3, 3.0, 0.45, size=13, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(s, "·  " + item, x+0.15, 2.88 + j*0.65, 2.7, 0.55,
                 size=10, color=C_MUTED)

# ─────────────────────────────────────────
# SLIDE 7 — Agile Process & Team
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)
slide_label(s, "AGILE PROCESS & TEAM")
title_block(s, "Scrum-Based Development — 3 Sprints",
            "Daily standups, weekly sprint reviews, and a full product backlog with 10 items.")

# Team cards
team = [
    ("Risly", "Product Owner",
     "Defined product vision, validated acceptance\ncriteria, wrote executive summary, tested UX."),
    ("Sultan", "Scrum Master + Backend",
     "Led sprints, fixed blockers, built auth system,\nhardened security, managed deployment."),
    ("Fathir", "Developer — AI + Backend",
     "Collected dataset, trained YOLO11, built FastAPI\nendpoints, OpenCV pipeline, WebSocket."),
    ("Misha", "Developer — Frontend",
     "Built entire React dashboard, integrated all APIs,\ndesigned UI/UX, responsive layout."),
]
for i, (name, role, desc) in enumerate(team):
    x = 0.5 + i*3.2
    add_rect(s, x, 2.3, 3.0, 2.8, fill=C_CARD, line=C_BORDER, line_w=0.5)
    add_text(s, name, x+0.15, 2.42, 2.7, 0.42, size=16, bold=True, color=C_WHITE)
    add_text(s, role, x+0.15, 2.86, 2.7, 0.32, size=9.5,
             bold=True, color=C_ACCENT)
    add_rect(s, x+0.15, 3.22, 2.7, 0.015, fill=C_BORDER)
    add_text(s, desc, x+0.15, 3.3, 2.7, 1.5, size=9, color=C_MUTED)

# Sprint timeline
sprints = [
    ("Sprint 1", "Dataset & Foundation",
     "Repo setup, dataset collection\n1,200 annotated images"),
    ("Sprint 2", "Model & Streaming",
     "YOLO training, MJPEG stream\nLogin, video integration"),
    ("Sprint 3", "Alerts & Hardening",
     "WebSocket alerts, auth hardening\nBug fixes, analytics, demo prep"),
]
add_rect(s, 0.5, 5.42, 12.33, 0.045, fill=C_ACCENT)
for i, (sp, title, detail) in enumerate(sprints):
    x = 0.5 + i*4.28
    add_rect(s, x+1.9, 5.28, 0.12, 0.28, fill=C_ACCENT)
    add_text(s, sp, x, 5.62, 4.0, 0.3, size=10, bold=True,
             color=C_ACCENT)
    add_text(s, title, x, 5.94, 4.0, 0.3, size=11, bold=True, color=C_WHITE)
    add_text(s, detail, x, 6.28, 4.0, 0.65, size=9, color=C_MUTED)

# ─────────────────────────────────────────
# SLIDE 8 — Product Backlog Summary
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)
slide_label(s, "PRODUCT BACKLOG")
title_block(s, "10 Backlog Items — All Completed",
            "Ordered by priority. High items are foundational; medium/low items add value and polish.")

rows = [
    ("PB-01", "Object Detection Model Training",       "High",   "Done", C_SUCCESS),
    ("PB-02", "Multi-Zone Camera Architecture",        "High",   "Done", C_SUCCESS),
    ("PB-03", "Real-Time Video Stream + Bounding Box", "High",   "Done", C_SUCCESS),
    ("PB-04", "Detection Event Logging to Database",   "High",   "Done", C_SUCCESS),
    ("PB-05", "Role-Based Authentication & Access",    "High",   "Done", C_SUCCESS),
    ("PB-06", "Real-Time WebSocket Alert System",      "High",   "Done", C_SUCCESS),
    ("PB-07", "Analytics Dashboard & Risk Reports",    "Medium", "Done", C_SUCCESS),
    ("PB-08", "Persistent Token & Password Reset",     "Medium", "Done", C_SUCCESS),
    ("PB-09", "AI Performance Monitoring Page",        "Medium", "Done", C_SUCCESS),
    ("PB-10", "Settings, Preferences & User Mgmt",    "Low",    "Done", C_SUCCESS),
]
# Header
add_rect(s, 0.5, 2.2, 12.33, 0.38, fill=C_ACCENT)
for col_x, col_w, label in [(0.5, 0.8, "ID"), (1.35, 6.5, "Description"),
                              (7.9, 1.8, "Priority"), (9.75, 1.5, "Status")]:
    add_text(s, label, col_x+0.1, 2.22, col_w, 0.32,
             size=9.5, bold=True, color=C_WHITE)

for i, (pid, desc, pri, status, sc) in enumerate(rows):
    y = 2.62 + i*0.46
    bg_c = C_CARD if i%2==0 else C_BG
    add_rect(s, 0.5, y, 12.33, 0.44, fill=bg_c)
    add_text(s, pid,    0.6,  y+0.07, 0.7,  0.3, size=8.5, bold=True, color=C_ACCENT)
    add_text(s, desc,   1.35, y+0.07, 6.4,  0.3, size=9,   color=C_WHITE)
    pri_col = C_DANGER if pri=="High" else C_WARN if pri=="Medium" else C_MUTED
    add_text(s, pri,    7.9,  y+0.07, 1.7,  0.3, size=8.5, color=pri_col)
    add_text(s, status, 9.75, y+0.07, 1.4,  0.3, size=8.5, bold=True, color=sc)

# ─────────────────────────────────────────
# SLIDE 9 — Business Value
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)
slide_label(s, "BUSINESS VALUE")
title_block(s, "Why SmartGuard AI Matters",
            "Real operational impact for PT. Kawan Lama and the broader logistics industry.")

impacts = [
    ("Response Time",  "From hours\nto seconds",    C_DANGER,
     "Automated detection means threats are identified the moment\nthey appear, not when a guard happens to look at the right screen."),
    ("Cost Reduction",  "Fewer\nfalse alarms",      C_WARN,
     "Confidence thresholds and cooldown timers eliminate noise.\nOnly real detections generate alerts and log entries."),
    ("Scalability",    "Unlimited\ncamera zones",   C_ACCENT,
     "Adding a new zone requires only a database entry. The worker\nthread architecture handles any number of parallel streams."),
    ("Audit Trail",    "Full incident\nhistory",     C_SUCCESS,
     "Every detection is timestamped, geotagged to a zone, saved\nwith a snapshot image, and accessible in the Detection Logs page."),
]
for i, (title, val, col, body) in enumerate(impacts):
    row = i // 2
    c   = i % 2
    x = 0.5 + c*6.45
    y = 2.3 + row*2.3
    add_rect(s, x, y, 6.15, 2.1, fill=C_CARD, line=col, line_w=1.0)
    add_text(s, val,   x+0.2, y+0.12, 2.5, 0.75, size=20, bold=True, color=col)
    add_text(s, title, x+0.2, y+0.88, 5.7, 0.3,  size=10, bold=True, color=C_WHITE)
    add_text(s, body,  x+0.2, y+1.2,  5.7, 0.75, size=9,  color=C_MUTED)

# ─────────────────────────────────────────
# SLIDE 10 — Closing
# ─────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_BG)
add_rect(s, 0, 0, 0.08, 7.5, fill=C_ACCENT)
add_rect(s, 0, 6.8, 13.33, 0.7, fill=C_CARD)

add_text(s, "Thank You", 0.5, 1.5, 12, 1.0,
         size=52, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(s, "SmartGuard AI — Keeping warehouses safe, automatically.",
         0.5, 2.65, 12, 0.5, size=18, color=C_ACCENT,
         align=PP_ALIGN.CENTER, bold=True)
add_rect(s, 4.5, 3.35, 4.33, 0.045, fill=C_ACCENT)

links = [
    "GitHub   :  github.com/SultanZhalifa/smartwarehouse-ai",
    "Group     :  5  (Risly · Sultan · Fathir · Misha)",
    "Course   :  AI Open Innovation Challenge 2026",
]
for i, line in enumerate(links):
    add_text(s, line, 0.5, 3.7 + i*0.55, 12, 0.45,
             size=12, color=C_MUTED, align=PP_ALIGN.CENTER)

add_text(s, "Group 5  ·  President University  ·  2026",
         0.5, 6.85, 12, 0.4, size=10, color=C_MUTED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# Save
# ─────────────────────────────────────────
out = r"c:\Tugas + Hackathon\smartwarehouse-ai\docs\SmartGuard AI - Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
